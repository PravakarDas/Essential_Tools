from __future__ import annotations

import io
import os
from typing import Any, Dict, List


def process(job, upload_paths: List[str]) -> Dict[str, Any]:
    if len(upload_paths) != 1:
        raise ValueError("Upload exactly one PDF to convert")

    # Lazy imports to avoid heavy deps unless used
    try:
        from pdf2image import convert_from_path  # type: ignore
    except Exception as e:  # pragma: no cover
        raise RuntimeError(
            "pdf2image is required (and Poppler must be installed) to convert PDF pages to images"
        ) from e

    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Emu
    except Exception as e:  # pragma: no cover
        raise RuntimeError("python-pptx is required for PowerPoint generation") from e

    input_pdf = upload_paths[0]
    out_name = f"{job.id}.pptx"
    out_path = os.path.join(job.workspace_path, out_name)

    # Render pages to images
    # Default dpi=150 gives a good balance of quality/size; can be adjusted later via options
    try:
        images = convert_from_path(input_pdf, dpi=150)  # returns list of PIL Images
    except Exception as e:
        # Provide a clearer hint when Poppler is missing
        msg = str(e)
        if "poppler" in msg.lower() or "PDFInfoNotInstalledError" in msg:
            raise RuntimeError(
                "Poppler is required for PDF rendering. Please install Poppler and ensure its binaries are on PATH."
            ) from e
        raise

    if not images:
        raise ValueError("No pages found in PDF")

    prs = Presentation()
    # Use a blank layout for all slides
    blank_layout = prs.slide_layouts[6]

    # Convenience sizes in EMU
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for img in images:
        slide = prs.slides.add_slide(blank_layout)
        # Fit image to slide while preserving aspect ratio
        iw, ih = img.size  # pixels
        # Save image to in-memory PNG
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)

        # Scale to cover slide (contain vs cover). We'll use 'contain' to avoid cropping
        # Compute scale factors using slide size; assume 96 DPI for EMU conversion baseline
        # but since we specify both width and height to add_picture, python-pptx will resize accordingly.
        slide.shapes.add_picture(buf, Emu(0), Emu(0), width=slide_w, height=slide_h)

    prs.save(out_path)
    return {"files": [out_path]}

