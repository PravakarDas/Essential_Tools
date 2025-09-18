from __future__ import annotations

import os
import shutil
import subprocess
from typing import Any, Dict, List


def _which(cmd: str) -> str | None:
    return shutil.which(cmd)


def _libreoffice_exe() -> str | None:
    for name in ("soffice", "soffice.bin"):
        p = _which(name)
        if p:
            return p
    return None


def _convert_with_libreoffice(inp: str, out_dir: str) -> str:
    exe = _libreoffice_exe()
    if not exe:
        raise FileNotFoundError("LibreOffice (soffice) not found")
    cmd = [exe, "--headless", "--nologo", "--convert-to", "pdf", "--outdir", out_dir, inp]
    subprocess.run(cmd, check=True)
    base = os.path.splitext(os.path.basename(inp))[0] + ".pdf"
    return os.path.join(out_dir, base)


EMU_PER_PT = 12700.0


def _pptx_to_pdf_pure(src: str, out_path: str) -> None:
    # Render PPTX slides into a PDF using python-pptx (read) + PyMuPDF (draw)
    from pptx import Presentation  # type: ignore
    import fitz  # type: ignore

    prs = Presentation(src)
    slide_w_pt = float(prs.slide_width) / EMU_PER_PT
    slide_h_pt = float(prs.slide_height) / EMU_PER_PT
    doc = fitz.open()

    for slide in prs.slides:
        page = doc.new_page(width=slide_w_pt, height=slide_h_pt)
        for shape in slide.shapes:
            try:
                left = float(getattr(shape, "left", 0)) / EMU_PER_PT
                top = float(getattr(shape, "top", 0)) / EMU_PER_PT
                width = float(getattr(shape, "width", 0)) / EMU_PER_PT
                height = float(getattr(shape, "height", 0)) / EMU_PER_PT
            except Exception:
                left = top = 0.0
                width = height = 0.0

            rect = fitz.Rect(left, top, left + max(width, 0.1), top + max(height, 0.1))

            # Pictures
            if hasattr(shape, "image"):
                try:
                    blob = shape.image.blob  # type: ignore[attr-defined]
                    page.insert_image(rect, stream=blob)
                    continue
                except Exception:
                    pass

            # Text
            if getattr(shape, "has_text_frame", False):
                try:
                    tf = shape.text_frame
                    # Gather plain text with line breaks
                    lines: list[str] = []
                    for p in tf.paragraphs:
                        runs = [run.text or "" for run in getattr(p, "runs", [])]
                        lines.append("".join(runs))
                    text = "\n".join(lines)
                    # Use insert_textbox with safe args; fall back to insert_text if it fails
                    try:
                        page.insert_textbox(rect, text, fontname="helv", fontsize=12, align=0)
                    except Exception:
                        # Simple top-left text draw as a fallback
                        y = rect.y0
                        lh = 14
                        for ln in text.splitlines() or [text]:
                            page.insert_text((rect.x0, y), ln, fontname="helv", fontsize=12)
                            y += lh
                except Exception:
                    pass

    doc.save(out_path)
    doc.close()


def process(job, upload_paths: List[str]) -> Dict[str, Any]:
    if len(upload_paths) != 1:
        raise ValueError("Upload exactly one PowerPoint file to convert")
    src = upload_paths[0]
    out_name = f"{job.id}.pdf"
    out_path = os.path.join(job.workspace_path, out_name)

    ext = os.path.splitext(src)[1].lower()

    # Try LibreOffice first when available (best fidelity for PPT/PPTX)
    try:
        if _libreoffice_exe():
            produced = _convert_with_libreoffice(src, job.workspace_path)
            if os.path.abspath(produced) != os.path.abspath(out_path):
                os.replace(produced, out_path)
            return {"files": [out_path]}
    except Exception:
        pass

    # Pure-Python fallback for PPTX using installed libraries (python-pptx + PyMuPDF)
    if ext == ".pptx":
        _pptx_to_pdf_pure(src, out_path)
        return {"files": [out_path]}

    # PPT (legacy) requires LibreOffice
    raise RuntimeError(
        "This conversion supports PPTX without extra libraries. For PPT, please install LibreOffice or save as PPTX."
    )
